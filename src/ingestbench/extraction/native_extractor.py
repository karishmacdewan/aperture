"""Native extractor: pdfplumber for PDF_NATIVE, python-docx for DOCX, python-pptx for PPTX.

No OCR, no credentials required -- if the file has a text layer we extract it directly.
"""

from __future__ import annotations

import time
from pathlib import Path

from ingestbench.core.interfaces import Extractor
from ingestbench.core.models import ExtractionResult, FileType
from ingestbench.core.registry import registry


class NativeExtractor(Extractor):
    name = "native"
    supported_types = [FileType.PDF_NATIVE, FileType.DOCX, FileType.PPTX, FileType.TEXT_MARKDOWN]
    requires_credentials = False

    def extract(self, file_path: Path) -> ExtractionResult:
        from ingestbench.extraction.file_types import detect_file_type

        t0 = time.perf_counter()
        file_type = detect_file_type(file_path)
        try:
            if file_type == FileType.PDF_NATIVE:
                result = self._extract_pdf(file_path)
            elif file_type == FileType.DOCX:
                result = self._extract_docx(file_path)
            elif file_type == FileType.PPTX:
                result = self._extract_pptx(file_path)
            elif file_type == FileType.TEXT_MARKDOWN:
                result = self._extract_text(file_path)
            else:
                result = ExtractionResult(
                    extractor_name=self.name,
                    file_name=file_path.name,
                    file_type=file_type,
                    success=False,
                    quality_flags=[f"unsupported_file_type:{file_type}"],
                )
        except Exception as exc:
            result = ExtractionResult(
                extractor_name=self.name,
                file_name=file_path.name,
                file_type=file_type,
                success=False,
                quality_flags=[f"extraction_error:{exc}"],
            )

        result.extraction_time_s = round(time.perf_counter() - t0, 4)
        return result

    def _extract_pdf(self, file_path: Path) -> ExtractionResult:
        import pdfplumber

        pages_text: list[str] = []
        tables_detected = 0
        failed_pages: list[int] = []

        with pdfplumber.open(str(file_path)) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    text = page.extract_text() or ""
                    pages_text.append(text)
                    tables_detected += len(page.extract_tables() or [])
                except Exception:
                    failed_pages.append(i + 1)
                    pages_text.append("")

        raw_text = "\n\n".join(t for t in pages_text if t.strip())
        return ExtractionResult(
            extractor_name=self.name,
            file_name=file_path.name,
            file_type=FileType.PDF_NATIVE,
            success=True,
            char_count=len(raw_text),
            page_count=page_count,
            tables_detected=tables_detected,
            failed_pages=failed_pages,
            raw_text=raw_text,
            config_used={"library": "pdfplumber"},
        )

    def _extract_docx(self, file_path: Path) -> ExtractionResult:
        from docx import Document

        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_detected = len(doc.tables)
        # Also pull text from table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)

        raw_text = "\n\n".join(paragraphs)
        return ExtractionResult(
            extractor_name=self.name,
            file_name=file_path.name,
            file_type=FileType.DOCX,
            success=True,
            char_count=len(raw_text),
            page_count=0,  # DOCX has no fixed page concept at parse time
            tables_detected=tables_detected,
            raw_text=raw_text,
            config_used={"library": "python-docx"},
        )

    def _extract_pptx(self, file_path: Path) -> ExtractionResult:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        slide_texts: list[str] = []
        images_detected = 0

        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            parts.append(text)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    images_detected += 1
            slide_texts.append("\n".join(parts))

        raw_text = "\n\n".join(t for t in slide_texts if t.strip())
        return ExtractionResult(
            extractor_name=self.name,
            file_name=file_path.name,
            file_type=FileType.PPTX,
            success=True,
            char_count=len(raw_text),
            page_count=len(prs.slides),
            images_detected=images_detected,
            raw_text=raw_text,
            config_used={"library": "python-pptx"},
        )

    def _extract_text(self, file_path: Path) -> ExtractionResult:
        raw_text = file_path.read_text(encoding="utf-8", errors="replace")
        return ExtractionResult(
            extractor_name=self.name,
            file_name=file_path.name,
            file_type=FileType.TEXT_MARKDOWN,
            success=True,
            char_count=len(raw_text),
            raw_text=raw_text,
            config_used={"library": "builtin"},
        )


registry.register_extractor(NativeExtractor())
