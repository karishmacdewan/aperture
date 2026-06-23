"""Synthetic benchmark document generator.

Deterministically produces a small, reproducible document corpus across four
categories, each exercising a different part of the extraction layer (see
ARCHITECTURE.md, Section 12). Each category's output directory also gets a
``ground_truth.json`` manifest recording exactly what was generated (heading
count, table shape, page count, the exact text rendered, etc.), so the
extraction comparison in later phases can report a real accuracy figure
instead of a self-reported confidence score.

Usage:
    python sample_docs/generate_synthetic_docs.py [--out-dir DIR]

Requires the "sample-docs" extra:
    pip install -e ".[sample-docs]"
"""

from __future__ import annotations

import argparse
import json
import tempfile
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table as RLTable,
    TableStyle,
    Image as RLImage,
)
from pypdf import PdfReader

try:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except ImportError:  # matplotlib is part of the sample-docs extra; fail loudly later if missing
    plt = None


# ---------------------------------------------------------------------------
# Shared content fixtures (kept here, not randomly generated, so the same
# logical content can be rendered into multiple formats and the ground truth
# is known by construction).
# ---------------------------------------------------------------------------

STRUCTURED_TITLE = "Acme Corp Data Retention Policy"

STRUCTURED_SECTIONS: list[dict[str, Any]] = [
    {
        "heading": "1. Introduction",
        "body": (
            "This policy defines how Acme Corp classifies, retains, and "
            "disposes of company records. It applies to all employees, "
            "contractors, and systems that create or store company data."
        ),
    },
    {
        "heading": "2. Scope",
        "body": (
            "This policy covers structured and unstructured records held in "
            "company systems, including customer records, email logs, and "
            "financial statements. Personal devices are out of scope."
        ),
    },
    {
        "heading": "3. Retention Requirements",
        "body": (
            "Each data category below has a minimum retention period and a "
            "designated owner responsible for enforcement and disposal."
        ),
    },
    {
        "heading": "4. Appendix: Exceptions",
        "body": (
            "Records subject to active legal hold are retained until the "
            "hold is released, regardless of the periods listed above."
        ),
    },
]

STRUCTURED_TABLE_HEADER = ["Data Category", "Retention Period", "Owner"]
STRUCTURED_TABLE_ROWS = [
    ["Customer Records", "7 years", "Legal"],
    ["Email Logs", "1 year", "IT Security"],
    ["Financial Statements", "10 years", "Finance"],
]

SCANNED_PAGES: list[dict[str, Any]] = [
    {
        "heading": "MEMORANDUM",
        "lines": [
            "To: All Regional Managers",
            "From: Operations Planning",
            "Re: Warehouse Inventory Audit Schedule",
            "",
            "All regional warehouses must complete a full inventory audit",
            "before the end of the current quarter. Audit teams should",
            "coordinate with the regional manager to schedule downtime.",
        ],
    },
    {
        "heading": "Page 2 - Continued",
        "lines": [
            "Audit reports are due five business days after the audit",
            "is completed. Late submissions require a written explanation",
            "and a copy to the regional director.",
            "",
            "Questions should be directed to Operations Planning.",
        ],
    },
]

PPTX_SLIDES: list[dict[str, Any]] = [
    {
        "type": "title",
        "title": "Q3 Infrastructure Review",
        "subtitle": "Internal Engineering Update",
    },
    {
        "type": "bullets",
        "title": "Highlights",
        "bullets": [
            "Migrated 60% of services to the new cluster",
            "Reduced p95 latency by 18%",
            "3 incidents resolved within SLA",
        ],
    },
    {
        "type": "chart",
        "title": "Uptime by Quarter",
        "bullets": ["Uptime trended upward across all three measured quarters"],
    },
    {
        "type": "bullets",
        "title": "Next Steps",
        "bullets": [
            "Complete migration by Q4",
            "Run load test in staging",
            "Publish updated runbook",
        ],
    },
]

IMAGE_HEAVY_DOCS: list[dict[str, Any]] = [
    {
        "filename": "screenshot_1.png",
        "lines": [
            "build-server-04 ~ % deploy status",
            "service: ingestion-api      state: RUNNING   uptime: 14d 02h",
            "service: embedding-worker   state: RUNNING   uptime: 14d 02h",
            "service: report-generator   state: DEGRADED  uptime: 00d 03h",
            "",
            "WARN report-generator: queue depth above threshold (412 > 250)",
            "INFO autoscaler: spinning up 1 additional worker",
        ],
        "diagram_label": "ingestion -> queue -> worker",
    },
    {
        "filename": "screenshot_2.png",
        "lines": [
            "Pipeline Diagram (manual export)",
            "",
            "Documents -> Extraction -> Chunking -> Embedding -> Vector DB",
            "",
            "Each stage logs duration_s and status to runs/<run_id>/run.log",
        ],
        "diagram_label": "stage 1 -> stage 2 -> stage 3",
    },
]


def _load_font(size: int) -> ImageFont.FreeTypeFont:
    """Best-effort load of a real TrueType font so rendered text is legible
    and OCR-able; falls back to PIL's bitmap default font if none is found.
    """
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


# ---------------------------------------------------------------------------
# Category 1: Structured documents (DOCX + machine-readable PDF)
# ---------------------------------------------------------------------------

def generate_structured_docs(out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)

    # --- DOCX ---
    docx_path = out_dir / "sample_structured.docx"
    doc = DocxDocument()
    doc.add_heading(STRUCTURED_TITLE, level=0)
    for section in STRUCTURED_SECTIONS:
        doc.add_heading(section["heading"], level=1)
        doc.add_paragraph(section["body"])
        if section["heading"].startswith("3."):
            table = doc.add_table(rows=1, cols=len(STRUCTURED_TABLE_HEADER))
            table.style = "Light Grid Accent 1"
            hdr_cells = table.rows[0].cells
            for i, col_name in enumerate(STRUCTURED_TABLE_HEADER):
                hdr_cells[i].text = col_name
            for row in STRUCTURED_TABLE_ROWS:
                cells = table.add_row().cells
                for i, val in enumerate(row):
                    cells[i].text = val
    doc.save(docx_path)

    # --- Machine-readable PDF (same logical content) ---
    pdf_path = out_dir / "sample_structured.pdf"
    styles = getSampleStyleSheet()
    h0 = ParagraphStyle("h0", parent=styles["Title"], fontSize=18, spaceAfter=16)
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontSize=13, spaceAfter=8, spaceBefore=12)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=10, spaceAfter=8)

    flow = [Paragraph(STRUCTURED_TITLE, h0)]
    for section in STRUCTURED_SECTIONS:
        flow.append(Paragraph(section["heading"], h1))
        flow.append(Paragraph(section["body"], body))
        if section["heading"].startswith("3."):
            table_data = [STRUCTURED_TABLE_HEADER] + STRUCTURED_TABLE_ROWS
            t = RLTable(table_data, hAlign="LEFT")
            t.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2F5496")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            flow.append(Spacer(1, 4))
            flow.append(t)
            flow.append(Spacer(1, 8))
    SimpleDocTemplate(str(pdf_path), pagesize=LETTER).build(flow)

    actual_pages = len(PdfReader(str(pdf_path)).pages)
    full_text = STRUCTURED_TITLE + "\n" + "\n".join(
        f"{s['heading']}\n{s['body']}" for s in STRUCTURED_SECTIONS
    )

    ground_truth = {
        "category": "structured",
        "files": [docx_path.name, pdf_path.name],
        "num_headings": 1 + len(STRUCTURED_SECTIONS),  # title + 4 section headings
        "headings": [STRUCTURED_TITLE] + [s["heading"] for s in STRUCTURED_SECTIONS],
        "num_tables": 1,
        "tables": [
            {
                "num_rows": 1 + len(STRUCTURED_TABLE_ROWS),  # header + data rows
                "num_cols": len(STRUCTURED_TABLE_HEADER),
                "header": STRUCTURED_TABLE_HEADER,
                "rows": STRUCTURED_TABLE_ROWS,
            }
        ],
        "num_pages_pdf": actual_pages,
        "ocr_required": False,
        "full_text": full_text,
        "notes": "Native text layer in both formats; no OCR should be triggered.",
    }
    (out_dir / "ground_truth.json").write_text(json.dumps(ground_truth, indent=2))
    return ground_truth


# ---------------------------------------------------------------------------
# Category 2: Scanned PDFs (text rendered to an image, no text layer)
# ---------------------------------------------------------------------------

def generate_scanned_pdf(out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    page_size = (1700, 2200)  # ~8.5x11in at 200dpi
    heading_font = _load_font(34)
    body_font = _load_font(24)

    pdf_path = out_dir / "sample_scanned.pdf"

    # Intermediate page images are scratch files -- written to a real tmp
    # directory (not the output mount) so they can be cleanly discarded once
    # the PDF is assembled.
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        page_image_paths: list[Path] = []
        for i, page in enumerate(SCANNED_PAGES, start=1):
            img = Image.new("L", page_size, color=255)
            draw = ImageDraw.Draw(img)
            y = 120
            draw.text((100, y), page["heading"], font=heading_font, fill=0)
            y += 70
            for line in page["lines"]:
                draw.text((100, y), line, font=body_font, fill=0)
                y += 42
            img_path = tmp_dir / f"page_{i}.png"
            img.save(img_path)
            page_image_paths.append(img_path)

        doc = SimpleDocTemplate(str(pdf_path), pagesize=LETTER)
        flow = []
        for i, img_path in enumerate(page_image_paths):
            flow.append(RLImage(str(img_path), width=468, height=606))
            if i < len(page_image_paths) - 1:
                from reportlab.platypus import PageBreak

                flow.append(PageBreak())
        doc.build(flow)

    # Confirm this is a true scan: no extractable text layer.
    extracted = "".join(p.extract_text() or "" for p in PdfReader(str(pdf_path)).pages)
    has_text_layer = len(extracted.strip()) > 0

    ground_truth = {
        "category": "scanned",
        "files": [pdf_path.name],
        "num_pages": len(SCANNED_PAGES),
        "ocr_required": True,
        "has_text_layer": has_text_layer,
        "pages": [
            {"page_number": i + 1, "heading": pg["heading"], "lines": pg["lines"]}
            for i, pg in enumerate(SCANNED_PAGES)
        ],
        "full_text": "\n".join(
            pg["heading"] + "\n" + "\n".join(pg["lines"]) for pg in SCANNED_PAGES
        ),
        "notes": (
            "Page content is rasterized text with no underlying text layer "
            "-- a true scan. Extraction should route to OCR (Tesseract "
            "default, Azure DI if configured)."
        ),
    }
    (out_dir / "ground_truth.json").write_text(json.dumps(ground_truth, indent=2))
    return ground_truth


# ---------------------------------------------------------------------------
# Category 3: PowerPoint (with one embedded chart image)
# ---------------------------------------------------------------------------

def _make_chart_image(path: Path) -> None:
    if plt is None:
        raise RuntimeError("matplotlib is required to generate the PPTX chart image")
    fig, ax = plt.subplots(figsize=(6, 3.5))
    quarters = ["Q1", "Q2", "Q3"]
    uptime = [99.2, 99.5, 99.8]
    ax.bar(quarters, uptime, color="#2F5496")
    ax.set_ylim(98, 100)
    ax.set_ylabel("Uptime %")
    ax.set_title("Uptime by Quarter")
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def generate_powerpoint(out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / "sample_deck.pptx"
    chart_slide_index = None

    # The chart image is a scratch file embedded into the .pptx -- written to
    # a real tmp directory (not the output mount) so it can be cleanly
    # discarded once embedded.
    with tempfile.TemporaryDirectory() as tmp:
        chart_path = Path(tmp) / "chart.png"
        _make_chart_image(chart_path)

        prs = Presentation()
        for i, slide_spec in enumerate(PPTX_SLIDES):
            if slide_spec["type"] == "title":
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_spec["title"]
                slide.placeholders[1].text = slide_spec["subtitle"]
            else:
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_spec["title"]
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = slide_spec["bullets"][0]
                for bullet in slide_spec["bullets"][1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = PptxPt(18)
                if slide_spec["type"] == "chart":
                    chart_slide_index = i
                    slide.shapes.add_picture(
                        str(chart_path), Inches(5.0), Inches(1.5), width=Inches(4.0)
                    )

        prs.save(pptx_path)

    ground_truth = {
        "category": "powerpoint",
        "files": [pptx_path.name],
        "num_slides": len(PPTX_SLIDES),
        "slides": [
            {
                "index": i,
                "title": s["title"],
                "bullets": s.get("bullets", []),
                "has_chart_image": s["type"] == "chart",
            }
            for i, s in enumerate(PPTX_SLIDES)
        ],
        "chart_slide_index": chart_slide_index,
        "num_embedded_images": 1,
        "full_text": "\n".join(
            s["title"] + ("\n" + s.get("subtitle", "") if "subtitle" in s else "")
            + "\n".join(s.get("bullets", []))
            for s in PPTX_SLIDES
        ),
        "notes": "One slide contains an embedded chart image (bar chart, no text alternative).",
    }
    (out_dir / "ground_truth.json").write_text(json.dumps(ground_truth, indent=2))
    return ground_truth


# ---------------------------------------------------------------------------
# Category 4: Image-heavy documents (dense text + a simple diagram)
# ---------------------------------------------------------------------------

def generate_image_heavy(out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    size = (1200, 800)
    text_font = _load_font(22)
    label_font = _load_font(18)

    entries = []
    for spec in IMAGE_HEAVY_DOCS:
        img = Image.new("RGB", size, color="white")
        draw = ImageDraw.Draw(img)

        y = 40
        for line in spec["lines"]:
            draw.text((40, y), line, font=text_font, fill="black")
            y += 32

        # Simple diagram: three boxes connected by arrows, with a label.
        box_w, box_h = 160, 70
        box_y = size[1] - 180
        xs = [60, 300, 540]
        for x in xs:
            draw.rectangle(
                [x, box_y, x + box_w, box_y + box_h], outline="black", width=3
            )
        for x in xs[:-1]:
            y_mid = box_y + box_h // 2
            draw.line(
                [(x + box_w, y_mid), (x + 240, y_mid)], fill="black", width=3
            )
            draw.polygon(
                [
                    (x + 240, y_mid),
                    (x + 230, y_mid - 8),
                    (x + 230, y_mid + 8),
                ],
                fill="black",
            )
        draw.text((60, box_y + box_h + 12), spec["diagram_label"], font=label_font, fill="black")

        img_path = out_dir / spec["filename"]
        img.save(img_path)
        entries.append(
            {
                "filename": spec["filename"],
                "lines": spec["lines"],
                "diagram_present": True,
                "diagram_label": spec["diagram_label"],
                "width": size[0],
                "height": size[1],
            }
        )

    ground_truth = {
        "category": "image_heavy",
        "files": [e["filename"] for e in entries],
        "num_images": len(entries),
        "images": entries,
        "full_text": "\n\n".join(
            "\n".join(e["lines"]) + "\n" + e["diagram_label"] for e in entries
        ),
        "ocr_required": True,
        "notes": (
            "Standalone images with dense text and a simple flow diagram. "
            "Native/Docling text extraction is expected to recover ~0 "
            "characters here -- this category specifically exercises the "
            "GPT-4o VLM extractor."
        ),
    }
    (out_dir / "ground_truth.json").write_text(json.dumps(ground_truth, indent=2))
    return ground_truth


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def generate_all(out_dir: Path) -> dict[str, Any]:
    summary = {
        "structured": generate_structured_docs(out_dir / "structured"),
        "scanned": generate_scanned_pdf(out_dir / "scanned"),
        "powerpoint": generate_powerpoint(out_dir / "powerpoint"),
        "image_heavy": generate_image_heavy(out_dir / "image_heavy"),
    }
    return summary


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate the synthetic benchmark corpus.")
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=Path(__file__).resolve().parent,
        help="Directory under which structured/scanned/powerpoint/image_heavy subfolders are created.",
    )
    args = parser.parse_args()

    summary = generate_all(args.out_dir)

    print("Synthetic benchmark corpus generated:")
    for category, gt in summary.items():
        files = ", ".join(gt["files"])
        print(f"  - {category}: {files}")
    print(f"\nGround-truth manifests written alongside each category's files.")


if __name__ == "__main__":
    main()
